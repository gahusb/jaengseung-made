#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
================================================================================
  PPT 제작 자동화 도구 v1.0
  Made by 쟁승메이드  |  jaengseung-made.com
================================================================================

  필요 패키지 설치:
    pip install python-pptx openpyxl

  사용법:
    1. 아래 ── 설정 ── 영역에서 옵션을 수정하세요.
    2. data.xlsx 파일을 준비하세요 (형식: A열=슬라이드 제목, B~열=불릿 내용).
       → 파일이 없으면 예시 데이터로 자동 실행됩니다.
    3. 터미널에서 실행:  python ppt_automation_v1.0.py
    4. 같은 폴더에 PPT 파일이 저장됩니다.

  지원 기능:
    - 표지 / 내용 / 마무리 슬라이드 자동 생성
    - 엑셀에서 데이터 읽어 슬라이드 일괄 생성
    - 16:9 비율, 맞은 고딕 폰트
    - 색상 테마 커스터마이징 가능
    - 슬라이드 번호 자동 추가

  맞춤 개발이 필요하다면: jaengseung-made.com/freelance
================================================================================
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import openpyxl
from datetime import datetime
import logging
import sys
import os

# ── 설정 (이 부분을 수정하세요) ───────────────────────────────────────────────

DATA_FILE   = "data.xlsx"       # 입력 엑셀 파일 (없으면 예시 데이터 사용)
OUTPUT_FILE = f"발표자료_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"

# 표지 정보
TITLE_TEXT    = "발표 제목을 입력하세요"
SUBTITLE_TEXT = "부제목 또는 발표자 이름"
DATE_TEXT     = datetime.now().strftime("%Y년 %m월 %d일")
CONTACT_TEXT  = "jaengseung-made.com  |  문의: bgg8988@gmail.com"

# 슬라이드 크기 (16:9)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── 색상 테마 ─────────────────────────────────────────────────────────────────
# 원하는 색상으로 변경하세요 (RGB)
COLOR_PRIMARY   = RGBColor(0x1D, 0x4E, 0xD8)   # 파란색 (헤더, 강조)
COLOR_SECONDARY = RGBColor(0x0F, 0x17, 0x2A)   # 다크 네이비 (표지 배경)
COLOR_ACCENT    = RGBColor(0x60, 0xA5, 0xFA)   # 라이트 블루 (서브 강조)
COLOR_TEXT      = RGBColor(0x1E, 0x29, 0x3B)   # 진한 슬레이트 (본문)
COLOR_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BG        = RGBColor(0xF1, 0xF5, 0xF9)   # 연한 배경
COLOR_BULLET    = RGBColor(0x1D, 0x4E, 0xD8)   # 불릿 색상

FONT_NAME = "맑은 고딕"   # 한글 폰트 (시스템에 설치된 폰트명)

# ────────────────────────────────────────────────────────────────────────────

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[logging.StreamHandler(sys.stdout)],
)
logger = logging.getLogger(__name__)


# ── 헬퍼 함수 ─────────────────────────────────────────────────────────────────

def rgb(r: int, g: int, b: int) -> RGBColor:
    return RGBColor(r, g, b)


def add_rect(slide, left, top, width, height,
             fill: RGBColor | None = None,
             line: RGBColor | None = None):
    """사각형 도형 추가"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text: str,
             left, top, width, height,
             size: int = 18,
             bold: bool = False,
             color: RGBColor = COLOR_TEXT,
             align=PP_ALIGN.LEFT,
             italic: bool = False) -> None:
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size     = Pt(size)
    run.font.bold     = bold
    run.font.italic   = italic
    run.font.name     = FONT_NAME
    run.font.color.rgb = color


# ── 슬라이드 생성 함수 ────────────────────────────────────────────────────────

def create_title_slide(prs: Presentation, title: str, subtitle: str, date: str) -> None:
    """표지 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃

    # 배경
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COLOR_SECONDARY)

    # 왼쪽 강조 세로선
    add_rect(slide,
             left=Inches(1.2), top=Inches(2.2),
             width=Inches(0.06), height=Inches(3.0),
             fill=COLOR_ACCENT)

    # 타이틀
    add_text(slide, title,
             left=Inches(1.5), top=Inches(2.2),
             width=Inches(10.5), height=Inches(1.8),
             size=42, bold=True, color=COLOR_WHITE)

    # 서브타이틀
    add_text(slide, subtitle,
             left=Inches(1.5), top=Inches(4.1),
             width=Inches(10.5), height=Inches(0.9),
             size=22, color=COLOR_ACCENT, italic=True)

    # 구분선
    add_rect(slide,
             left=Inches(1.5), top=Inches(5.2),
             width=Inches(10.5), height=Inches(0.015),
             fill=rgb(0x1E, 0x40, 0xAF))

    # 날짜
    add_text(slide, date,
             left=Inches(1.5), top=Inches(5.4),
             width=Inches(6), height=Inches(0.6),
             size=13, color=rgb(0x94, 0xA3, 0xB8))

    logger.info("  📋 표지 슬라이드 생성")


def create_content_slide(prs: Presentation,
                         title: str,
                         bullets: list[str],
                         slide_num: int) -> None:
    """내용 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COLOR_BG)

    # 상단 헤더
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), fill=COLOR_PRIMARY)

    # 슬라이드 번호 (우상단)
    add_text(slide, f"{slide_num:02d}",
             left=Inches(11.8), top=Inches(0.18),
             width=Inches(1.3), height=Inches(0.85),
             size=30, bold=True, color=COLOR_ACCENT,
             align=PP_ALIGN.RIGHT)

    # 제목
    add_text(slide, title,
             left=Inches(0.7), top=Inches(0.22),
             width=Inches(10.8), height=Inches(0.8),
             size=24, bold=True, color=COLOR_WHITE)

    # 흰색 콘텐츠 박스
    add_rect(slide,
             left=Inches(0.7), top=Inches(1.4),
             width=Inches(11.9), height=Inches(5.7),
             fill=COLOR_WHITE)

    # 불릿 포인트
    MAX_BULLETS = 8
    for i, bullet in enumerate(bullets[:MAX_BULLETS]):
        y = Inches(1.75 + i * 0.65)

        # 불릿 마크 (작은 사각형)
        add_rect(slide,
                 left=Inches(0.95), top=y + Inches(0.22),
                 width=Inches(0.12), height=Inches(0.12),
                 fill=COLOR_BULLET)

        # 불릿 텍스트
        add_text(slide, bullet,
                 left=Inches(1.25), top=y,
                 width=Inches(11.1), height=Inches(0.62),
                 size=16, color=COLOR_TEXT)

    # 하단 라인
    add_rect(slide,
             left=Inches(0.7), top=Inches(7.0),
             width=Inches(11.9), height=Inches(0.015),
             fill=COLOR_PRIMARY)

    logger.info(f"  📄 슬라이드 {slide_num} 생성: {title}")


def create_divider_slide(prs: Presentation, chapter: str, label: str = "") -> None:
    """챕터 구분 슬라이드 (선택사항)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COLOR_PRIMARY)
    add_rect(slide,
             left=0, top=Inches(3.5),
             width=SLIDE_W, height=Inches(0.03),
             fill=COLOR_ACCENT)

    if label:
        add_text(slide, label,
                 left=Inches(0), top=Inches(2.5),
                 width=SLIDE_W, height=Inches(0.6),
                 size=14, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)

    add_text(slide, chapter,
             left=Inches(0), top=Inches(3.0),
             width=SLIDE_W, height=Inches(1.2),
             size=38, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)


def create_closing_slide(prs: Presentation, contact: str) -> None:
    """마무리 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COLOR_SECONDARY)
    add_rect(slide,
             left=0, top=Inches(3.55),
             width=SLIDE_W, height=Inches(0.03),
             fill=COLOR_ACCENT)

    add_text(slide, "감사합니다",
             left=Inches(0), top=Inches(2.5),
             width=SLIDE_W, height=Inches(1.0),
             size=52, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, contact,
             left=Inches(0), top=Inches(3.9),
             width=SLIDE_W, height=Inches(0.7),
             size=16, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)

    logger.info("  🎬 마무리 슬라이드 생성")


# ── 데이터 로드 ───────────────────────────────────────────────────────────────

EXAMPLE_DATA = [
    {
        "title": "시장 현황 분석",
        "bullets": [
            "2024년 국내 시장 규모: 1조 2,800억 원 (전년비 +18.3%)",
            "상위 3개사 점유율 합계: 61.4% — 과점 구조 지속",
            "B2B 부문 성장률: B2C 대비 2.3배 높은 성장세",
            "주요 고객층: 중소기업 및 스타트업 비중 확대 중",
        ],
    },
    {
        "title": "핵심 문제 정의",
        "bullets": [
            "운영 비용 연평균 15.2% 상승 → 수익성 압박",
            "고객 이탈률 22% — 업계 평균(14%)보다 높음",
            "내부 반복 업무에 월 평균 220시간 소요 (비효율)",
            "경쟁사 대비 디지털 전환 12개월 지연 상태",
        ],
    },
    {
        "title": "제안 솔루션",
        "bullets": [
            "Phase 1: 업무 자동화 도입 — 반복 업무 70% 자동화",
            "Phase 2: 고객 데이터 플랫폼(CDP) 구축 — 이탈 예측",
            "Phase 3: 실시간 대시보드 도입 — 의사결정 속도 향상",
            "예상 ROI: 투자 대비 320% (12개월 기준)",
        ],
    },
    {
        "title": "추진 일정 및 기대 효과",
        "bullets": [
            "1단계 (1~2개월): 현황 분석 및 시스템 설계",
            "2단계 (3~4개월): 파일럿 운영 및 피드백 수집",
            "3단계 (5~6개월): 전사 확대 및 고도화",
            "연간 비용 절감 목표: 약 4.5억 원",
            "고객 이탈률 목표: 22% → 12% 이하",
        ],
    },
]


def load_from_excel(filepath: str) -> list[dict]:
    """엑셀 파일에서 슬라이드 데이터 로드 (A열=제목, B열~=불릿)"""
    if not os.path.exists(filepath):
        logger.warning(f"⚠️  '{filepath}' 파일 없음 → 예시 데이터로 실행합니다.")
        return EXAMPLE_DATA

    wb = openpyxl.load_workbook(filepath)
    ws = wb.active
    slides = []

    for row in ws.iter_rows(min_row=2, values_only=True):
        title = str(row[0] or "").strip()
        if not title:
            continue
        bullets = [str(c).strip() for c in row[1:] if c and str(c).strip()]
        slides.append({"title": title, "bullets": bullets})

    logger.info(f"엑셀 로드 완료: {len(slides)}개 슬라이드 데이터")
    return slides


# ── 메인 ──────────────────────────────────────────────────────────────────────

def main():
    logger.info("=" * 60)
    logger.info("  PPT 제작 자동화 도구 v1.0  |  쟁승메이드")
    logger.info("=" * 60)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # 표지
    create_title_slide(prs, TITLE_TEXT, SUBTITLE_TEXT, DATE_TEXT)

    # 내용 슬라이드
    slides_data = load_from_excel(DATA_FILE)
    for i, slide in enumerate(slides_data, start=1):
        create_content_slide(prs, slide["title"], slide["bullets"], slide_num=i)

    # 마무리
    create_closing_slide(prs, CONTACT_TEXT)

    prs.save(OUTPUT_FILE)
    total = len(slides_data) + 2  # 표지 + 내용 + 마무리
    logger.info(f"\n✅ 저장 완료: {OUTPUT_FILE}  ({total}슬라이드)")
    logger.info("\n맞춤 PPT 자동화가 필요하다면: jaengseung-made.com")


if __name__ == "__main__":
    main()
